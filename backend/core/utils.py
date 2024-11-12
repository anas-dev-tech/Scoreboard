# utils.py
import pandas as pd
from django.core.exceptions import FieldDoesNotExist
import pdfplumber
from pptx import Presentation



class ImportData(object):
    def __init__(self, file):
        self.file = file


    def excel_to_model_instances(self, model_class, columns):
        """
        Converts an Excel file to a list of model instances.

        Parameters:
            file (UploadedFile): The Excel file to read.
            model_class (Model): The Django model class to create instances of.
            columns (dict): A dictionary mapping Excel column names to model field names.
                            Example: {'Excel Column 1': 'model_field_1', ...}

        Returns:
            list: A list of model instances with populated fields.
        """
        try:
            # Read the Excel file into a DataFrame
            df = pd.read_excel(self.file)
            
            # Verify that all specified columns exist in the model
            for excel_col, model_field in columns.items():
                if excel_col not in df.columns:
                    raise ValueError(f"Excel file is missing required column: '{excel_col}'")
                if not hasattr(model_class, model_field):
                    raise FieldDoesNotExist(f"Model '{model_class.__name__}' has no field '{model_field}'")
            
            # Create a list to store model instances
            instances = []

            # Iterate over DataFrame rows and create model instances
            for _, row in df.iterrows():
                # Create instance of the model class
                instance = model_class()

                # Populate the fields as per columns mapping
                for excel_col, model_field in columns.items():
                    setattr(instance, model_field, row[excel_col])

                # Add instance to the list
                instances.append(instance)

            return instances
        
        except Exception as e:
            raise ValueError(f"Error processing file: {e}")

    def pdf_to_text(self):
        text = ""
        with pdfplumber.open(self.file) as pdf:
            for page in pdf.pages:
                text += page.extract_text()
        return text

    
    def pptx_to_text(self):
        text = ""
        presentation = Presentation(self.file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text